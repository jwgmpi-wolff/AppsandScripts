import json
import logging
import os
import shutil
import sqlite3
import subprocess
import tempfile
import time
from collections import Counter
from datetime import datetime, timedelta, timezone
from html import escape
from pathlib import Path
from typing import Any
from urllib.error import HTTPError, URLError
from urllib.parse import quote, urlencode
from urllib.request import Request, urlopen

import azure.functions as func

app = func.FunctionApp(http_auth_level=func.AuthLevel.FUNCTION)

HTTP_ROUTE_AUTH_LEVEL = (
    func.AuthLevel.ANONYMOUS
    if (os.getenv("AZURE_FUNCTIONS_ENVIRONMENT") or os.getenv("FUNCTIONS_ENVIRONMENT") or "").strip().lower() == "development"
    else func.AuthLevel.FUNCTION
)

_artifact_dir_override = (os.getenv("FINOPS_ARTIFACT_DIR") or "").strip()
ARTIFACT_DIR = Path(_artifact_dir_override) if _artifact_dir_override else (Path(tempfile.gettempdir()) / "finops-artifacts")
try:
    ARTIFACT_DIR.mkdir(parents=True, exist_ok=True)
except OSError:
    # Fall back to temp storage if the configured path is unavailable on the host.
    ARTIFACT_DIR = Path(tempfile.gettempdir()) / "finops-artifacts"
    ARTIFACT_DIR.mkdir(parents=True, exist_ok=True)
HOURS_PER_YEAR = 24 * 365
RETAIL_CACHE_TTL_HOURS = 24
_report_db_override = (os.getenv("FINOPS_REPORT_DB_PATH") or "").strip()
REPORT_DB_PATH = Path(_report_db_override) if _report_db_override else (ARTIFACT_DIR / "finops_reports.db")
REPORT_DB_PATH.parent.mkdir(parents=True, exist_ok=True)
_CURRENT_SUBSCRIPTION_ID = (os.getenv("AZURE_SUBSCRIPTION_ID") or "").strip()
MSAL_PUBLIC_CLIENT_ID = (os.getenv("FINOPS_MSAL_CLIENT_ID") or "").strip()


def _subscription_id_from_owner_name() -> str:
    owner = (os.getenv("WEBSITE_OWNER_NAME") or "").strip()
    if not owner:
        return ""
    candidate = owner.split("+", 1)[0].strip()
    if len(candidate) == 36 and candidate.count("-") == 4:
        return candidate
    return ""


def _resolve_az_executable() -> str:
    for candidate in ("az", "az.cmd", "az.exe"):
        resolved = shutil.which(candidate)
        if resolved:
            return resolved

    if os.name == "nt":
        fallback_paths = [
            r"C:\\Program Files (x86)\\Microsoft SDKs\\Azure\\CLI2\\wbin\\az.cmd",
            r"C:\\Program Files\\Microsoft SDKs\\Azure\\CLI2\\wbin\\az.cmd",
        ]
        for path in fallback_paths:
            if os.path.exists(path):
                return path

    raise RuntimeError(
        "Azure CLI executable was not found in PATH for the Functions worker. "
        "Install Azure CLI and ensure az/az.cmd is available to the host process."
    )


def run_az(*args: str, arm_access_token: str | None = None) -> Any:
    try:
        az_executable = _resolve_az_executable()
    except RuntimeError:
        return _run_az_fallback(*args, arm_access_token=arm_access_token)

    completed = subprocess.run([az_executable, *args], capture_output=True, text=True)
    if completed.returncode != 0:
        raise RuntimeError(completed.stderr.strip() or completed.stdout.strip() or "Azure CLI command failed")
    output = completed.stdout.strip()
    if not output:
        return []
    try:
        return json.loads(output)
    except json.JSONDecodeError:
        return output


def _get_arm_access_token(override_token: str | None = None) -> str:
    if override_token:
        return override_token

    resource = "https://management.azure.com/"
    identity_endpoint = os.getenv("IDENTITY_ENDPOINT")
    identity_header = os.getenv("IDENTITY_HEADER")
    token_error: Exception | None = None
    max_retries = 3

    for attempt in range(max_retries):
        try:
            if identity_endpoint and identity_header:
                token_url = f"{identity_endpoint}?api-version=2019-08-01&resource={quote(resource, safe='')}"
                request = Request(token_url, headers={"X-IDENTITY-HEADER": identity_header, "Metadata": "true"})
            else:
                token_url = f"http://169.254.169.254/metadata/identity/oauth2/token?api-version=2018-02-01&resource={quote(resource, safe='')}"
                request = Request(token_url, headers={"Metadata": "true"})

            with urlopen(request, timeout=30) as response:
                payload = json.loads(response.read().decode("utf-8"))
            token = payload.get("access_token")
            if token:
                return str(token)
            raise RuntimeError("Managed identity token response did not include access_token")
        except (HTTPError, URLError, TimeoutError, OSError, RuntimeError) as exc:
            token_error = exc
            if attempt < max_retries - 1:
                backoff_seconds = 2 ** attempt
                logging.warning(
                    "Managed identity token fetch failed (attempt %d/%d). Retrying in %ds: %s",
                    attempt + 1,
                    max_retries,
                    backoff_seconds,
                    exc,
                )
                time.sleep(backoff_seconds)

    raise RuntimeError(f"Managed identity token acquisition failed after {max_retries} attempts: {token_error}")


def _cleanup_old_artifacts(max_age_hours: int = 168) -> None:
    cutoff = datetime.now(timezone.utc) - timedelta(hours=max_age_hours)
    for artifact in ARTIFACT_DIR.glob("finops_*"):
        try:
            if not artifact.is_file():
                continue
            mtime = datetime.fromtimestamp(artifact.stat().st_mtime, tz=timezone.utc)
            if mtime < cutoff:
                artifact.unlink()
                logging.info("Cleaned up old artifact: %s", artifact.name)
        except Exception as exc:
            logging.warning("Failed to clean up artifact %s: %s", artifact.name, exc)


def _arm_get_paged(path: str, params: dict[str, str], arm_access_token: str | None = None) -> list[dict[str, Any]]:
    token = _get_arm_access_token(override_token=arm_access_token)
    base = "https://management.azure.com"
    next_url = f"{base}{path}?{urlencode(params)}"
    items: list[dict[str, Any]] = []

    while next_url:
        request = Request(next_url, headers={"Authorization": f"Bearer {token}"})
        with urlopen(request, timeout=60) as response:
            payload = json.loads(response.read().decode("utf-8"))
        page_items = payload.get("value") if isinstance(payload, dict) else None
        if isinstance(page_items, list):
            items.extend(page_items)
        next_link = payload.get("nextLink") if isinstance(payload, dict) else None
        next_url = str(next_link) if next_link else ""

    return items


def _arm_get(path: str, api_version: str, arm_access_token: str | None = None) -> dict[str, Any]:
    token = _get_arm_access_token(override_token=arm_access_token)
    url = f"https://management.azure.com{path}?api-version={quote(api_version, safe='')}"
    request = Request(url, headers={"Authorization": f"Bearer {token}"})
    with urlopen(request, timeout=60) as response:
        payload = json.loads(response.read().decode("utf-8"))
    return payload if isinstance(payload, dict) else {}


def _run_az_fallback(*args: str, arm_access_token: str | None = None) -> Any:
    global _CURRENT_SUBSCRIPTION_ID
    if len(args) >= 2 and args[0] == "account" and args[1] == "list":
        subscriptions = _arm_get_paged("/subscriptions", {"api-version": "2020-01-01"}, arm_access_token=arm_access_token)
        return [
            {
                "id": item.get("subscriptionId"),
                "name": item.get("displayName"),
                "tenantId": item.get("tenantId"),
            }
            for item in subscriptions
        ]

    if len(args) >= 4 and args[0] == "account" and args[1] == "set" and args[2] == "--subscription":
        _CURRENT_SUBSCRIPTION_ID = args[3]
        return []

    sub_id = _CURRENT_SUBSCRIPTION_ID or (os.getenv("AZURE_SUBSCRIPTION_ID") or "").strip()
    if not sub_id:
        raise RuntimeError("No subscription id is available for non-CLI execution")

    if len(args) >= 2 and args[0] == "resource" and args[1] == "list":
        resources = _arm_get_paged(f"/subscriptions/{sub_id}/resources", {"api-version": "2021-04-01"}, arm_access_token=arm_access_token)
        output: list[dict[str, Any]] = []
        for item in resources:
            output.append(
                {
                    "id": item.get("id"),
                    "name": item.get("name"),
                    "type": item.get("type"),
                    "location": item.get("location"),
                    "resourceGroup": item.get("resourceGroup"),
                    "subscriptionId": sub_id,
                }
            )
        return output

    if len(args) >= 3 and args[0] == "consumption" and args[1] == "usage" and args[2] == "list":
        logging.info("Azure CLI not available in host; cost usage collection is skipped in fallback mode.")
        return []

    if len(args) >= 3 and args[0] == "advisor" and args[1] == "recommendation" and args[2] == "list":
        recommendations = _arm_get_paged(
            f"/subscriptions/{sub_id}/providers/Microsoft.Advisor/recommendations",
            {"api-version": "2023-01-01"},
            arm_access_token=arm_access_token,
        )
        output = []
        for item in recommendations:
            props = item.get("properties") if isinstance(item.get("properties"), dict) else {}
            short_desc = props.get("shortDescription") if isinstance(props.get("shortDescription"), dict) else {}
            output.append(
                {
                    "id": item.get("id"),
                    "name": item.get("name"),
                    "category": props.get("category"),
                    "impact": props.get("impact"),
                    "severity": props.get("risk"),
                    "description": short_desc.get("problem"),
                    "solution": short_desc.get("solution"),
                    "resourceId": props.get("resourceMetadata", {}).get("resourceId") if isinstance(props.get("resourceMetadata"), dict) else None,
                    "annualSavingsAmount": props.get("annualSavingsAmount"),
                    "savingsAmount": props.get("savingsAmount"),
                    "extendedProperties": props.get("extendedProperties") if isinstance(props.get("extendedProperties"), dict) else {},
                }
            )
        return output

    if len(args) >= 5 and args[0] == "vm" and args[1] == "show" and args[2] == "--ids":
        vm_id = args[3]
        vm = _arm_get(vm_id, "2023-09-01", arm_access_token=arm_access_token)
        hardware = vm.get("properties", {}).get("hardwareProfile") if isinstance(vm.get("properties"), dict) else {}
        return hardware.get("vmSize")

    raise RuntimeError("Azure CLI is not available and requested command is unsupported in fallback mode")


def _month_range() -> tuple[str, str]:
    now = datetime.now(timezone.utc)
    end = now.date()
    start = (now.replace(day=1) - timedelta(days=1)).replace(day=1)
    return start.isoformat(), end.isoformat()


def _to_float(value: Any) -> float | None:
    if value is None:
        return None
    if isinstance(value, (int, float)):
        return float(value)
    text = str(value).strip().replace(",", "")
    if not text:
        return None
    try:
        return float(text)
    except ValueError:
        return None


def _extract_annual_savings(rec: dict[str, Any]) -> float | None:
    extended = rec.get("extendedProperties") if isinstance(rec.get("extendedProperties"), dict) else {}
    candidates = [
        rec.get("annualSavingsAmount"),
        rec.get("savingsAmount"),
        extended.get("annualSavingsAmount"),
        extended.get("AnnualSavingsAmount"),
        extended.get("annualSavings"),
        extended.get("AnnualSavings"),
        extended.get("savingsAmount"),
        extended.get("SavingsAmount"),
        extended.get("savings"),
        extended.get("Savings"),
    ]
    for candidate in candidates:
        parsed = _to_float(candidate)
        if parsed is not None and parsed >= 0:
            return parsed
    return None


def _add_commitment_savings_fields(recommendation: dict[str, Any]) -> None:
    annual_savings = _extract_annual_savings(recommendation)
    if annual_savings is None:
        recommendation["reservationSavings1Year"] = None
        recommendation["reservationSavings2Year"] = None
        recommendation["reservationSavings3Year"] = None
        recommendation["savingsPlanSavings1Year"] = None
        recommendation["savingsPlanSavings2Year"] = None
        recommendation["savingsPlanSavings3Year"] = None
        recommendation["savingsEstimateBasis"] = "advisor-annual-savings"
        recommendation["pricingCurrency"] = None
        recommendation["billingCurrency"] = None
        recommendation["currencyMatch"] = None
        recommendation["savingsCurrency"] = None
        recommendation["fxApplied"] = None
        recommendation["fxRateApplied"] = None
        recommendation["fxSource"] = None
        recommendation["estimatedAverageInstanceCount"] = None
        recommendation["observedUsageHoursInPeriod"] = None
        return

    recommendation["reservationSavings1Year"] = round(annual_savings, 2)
    recommendation["reservationSavings2Year"] = round(annual_savings * 2, 2)
    recommendation["reservationSavings3Year"] = round(annual_savings * 3, 2)
    recommendation["savingsPlanSavings1Year"] = round(annual_savings, 2)
    recommendation["savingsPlanSavings2Year"] = round(annual_savings * 2, 2)
    recommendation["savingsPlanSavings3Year"] = round(annual_savings * 3, 2)
    recommendation["savingsEstimateBasis"] = "advisor-annual-savings"
    recommendation["pricingCurrency"] = None
    recommendation["billingCurrency"] = None
    recommendation["currencyMatch"] = None
    recommendation["savingsCurrency"] = None
    recommendation["fxApplied"] = None
    recommendation["fxRateApplied"] = None
    recommendation["fxSource"] = None
    recommendation["estimatedAverageInstanceCount"] = None
    recommendation["observedUsageHoursInPeriod"] = None


def _normalize_region_name(region: str | None) -> str:
    if not region:
        return ""
    return str(region).replace(" ", "").lower()


def _ensure_retail_cache_table(conn: sqlite3.Connection) -> None:
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS retail_price_cache (
            vmSize TEXT NOT NULL,
            armRegionName TEXT NOT NULL,
            fetchedAt TEXT NOT NULL,
            ratesJson TEXT NOT NULL,
            PRIMARY KEY (vmSize, armRegionName)
        )
        """
    )


def _get_cached_vm_commitment_rates(vm_size: str, arm_region_name: str) -> dict[str, Any] | None:
    if not vm_size or not arm_region_name:
        return None

    conn = sqlite3.connect(REPORT_DB_PATH)
    try:
        _ensure_retail_cache_table(conn)
        row = conn.execute(
            "SELECT fetchedAt, ratesJson FROM retail_price_cache WHERE vmSize = ? AND armRegionName = ?",
            (vm_size, arm_region_name),
        ).fetchone()
    finally:
        conn.close()

    if not row:
        return None

    fetched_at_text, rates_json = row
    try:
        fetched_at = datetime.fromisoformat(str(fetched_at_text))
    except ValueError:
        return None
    if fetched_at.tzinfo is None:
        fetched_at = fetched_at.replace(tzinfo=timezone.utc)

    if datetime.now(timezone.utc) - fetched_at > timedelta(hours=RETAIL_CACHE_TTL_HOURS):
        return None

    try:
        parsed = json.loads(str(rates_json))
        return parsed if isinstance(parsed, dict) else None
    except json.JSONDecodeError:
        return None


def _set_cached_vm_commitment_rates(vm_size: str, arm_region_name: str, rates: dict[str, Any]) -> None:
    if not vm_size or not arm_region_name:
        return

    conn = sqlite3.connect(REPORT_DB_PATH)
    try:
        _ensure_retail_cache_table(conn)
        conn.execute(
            """
            INSERT OR REPLACE INTO retail_price_cache (vmSize, armRegionName, fetchedAt, ratesJson)
            VALUES (?, ?, ?, ?)
            """,
            (
                vm_size,
                arm_region_name,
                datetime.now(timezone.utc).isoformat(),
                json.dumps(rates),
            ),
        )
        conn.commit()
    finally:
        conn.close()


def _retail_prices_query(filter_expression: str) -> list[dict[str, Any]]:
    encoded_filter = quote(filter_expression, safe="()'= ")
    url = f"https://prices.azure.com/api/retail/prices?$filter={encoded_filter}"
    items: list[dict[str, Any]] = []

    while url:
        with urlopen(url, timeout=20) as response:
            payload = json.loads(response.read().decode("utf-8"))
        page_items = payload.get("Items", [])
        if isinstance(page_items, list):
            items.extend(page_items)
        next_page = payload.get("NextPageLink")
        url = str(next_page) if next_page else ""

    return items


def _pick_min_unit_price(items: list[dict[str, Any]]) -> float | None:
    prices: list[float] = []
    for item in items:
        price = _to_float(item.get("unitPrice"))
        if price is None or price <= 0:
            continue
        if str(item.get("unitOfMeasure") or "").lower().find("hour") == -1:
            continue
        prices.append(price)
    if not prices:
        return None
    return min(prices)


def _pick_currency_code(*item_lists: list[dict[str, Any]]) -> str | None:
    for items in item_lists:
        for item in items:
            currency = item.get("currencyCode")
            if currency:
                return str(currency)
    return None


def _normalize_currency_code(value: str | None) -> str | None:
    if not value:
        return None
    normalized = str(value).strip().upper()
    return normalized or None


def _load_fx_configuration() -> dict[str, Any]:
    source = (os.getenv("FINOPS_FX_SOURCE") or "env").strip().lower()
    default_base = _normalize_currency_code(os.getenv("FINOPS_FX_BASE")) or "USD"
    rates_payload: Any = None
    source_label = "none"

    if source == "url":
        rates_url = (os.getenv("FINOPS_FX_RATES_URL") or "").strip()
        if rates_url:
            try:
                with urlopen(rates_url, timeout=15) as response:
                    rates_payload = json.loads(response.read().decode("utf-8"))
                source_label = f"url:{rates_url}"
            except Exception as exc:
                logging.warning("Unable to load FX rates from URL %s: %s", rates_url, exc)
    if rates_payload is None:
        raw_json = (os.getenv("FINOPS_FX_RATES_JSON") or "").strip()
        if raw_json:
            try:
                rates_payload = json.loads(raw_json)
                source_label = "env:FINOPS_FX_RATES_JSON"
            except Exception as exc:
                logging.warning("Unable to parse FINOPS_FX_RATES_JSON: %s", exc)

    base_currency = default_base
    rate_values: Any = None
    if isinstance(rates_payload, dict):
        if isinstance(rates_payload.get("rates"), dict):
            base_currency = _normalize_currency_code(rates_payload.get("base")) or default_base
            rate_values = rates_payload.get("rates")
        else:
            rate_values = rates_payload

    rates: dict[str, float] = {}
    if isinstance(rate_values, dict):
        for code, value in rate_values.items():
            normalized_code = _normalize_currency_code(str(code))
            parsed = _to_float(value)
            if normalized_code and parsed is not None and parsed > 0:
                rates[normalized_code] = parsed

    rates[base_currency] = 1.0
    return {
        "baseCurrency": base_currency,
        "rates": rates,
        "source": source_label,
    }


def _convert_currency_amount(amount: float | None, from_currency: str | None, to_currency: str | None, fx_config: dict[str, Any]) -> tuple[float | None, float | None]:
    if amount is None:
        return None, None
    from_code = _normalize_currency_code(from_currency)
    to_code = _normalize_currency_code(to_currency)
    if not from_code or not to_code:
        return None, None
    if from_code == to_code:
        return round(amount, 2), 1.0

    rates = fx_config.get("rates") if isinstance(fx_config.get("rates"), dict) else {}
    from_rate = _to_float(rates.get(from_code))
    to_rate = _to_float(rates.get(to_code))
    if from_rate is None or to_rate is None or from_rate <= 0 or to_rate <= 0:
        return None, None

    amount_in_base = amount / from_rate
    converted_amount = amount_in_base * to_rate
    fx_rate = converted_amount / amount if amount else None
    return round(converted_amount, 2), round(fx_rate, 8) if fx_rate is not None else None


def _build_vm_usage_profile(costs: list[dict[str, Any]], start_date: str, end_date: str) -> dict[str, dict[str, float]]:
    try:
        period_days = max(1, (datetime.fromisoformat(end_date).date() - datetime.fromisoformat(start_date).date()).days + 1)
    except ValueError:
        period_days = 30
    period_hours = float(period_days * 24)

    usage_by_resource: dict[str, float] = {}
    for cost in costs:
        resource_id = str(cost.get("resourceId") or "").strip().lower()
        if not resource_id:
            continue
        meter_category = str(cost.get("meterCategory") or "").lower()
        if "virtual machines" not in meter_category:
            continue
        usage_quantity = _to_float(cost.get("usageQuantity"))
        if usage_quantity is None or usage_quantity <= 0:
            continue
        usage_by_resource[resource_id] = usage_by_resource.get(resource_id, 0.0) + usage_quantity

    profiles: dict[str, dict[str, float]] = {}
    for resource_id, usage_hours in usage_by_resource.items():
        avg_instance_count = max(0.0, usage_hours / period_hours)
        annualized_hours = avg_instance_count * HOURS_PER_YEAR
        profiles[resource_id] = {
            "usageHoursInPeriod": round(usage_hours, 2),
            "averageInstanceCount": round(avg_instance_count, 4),
            "annualizedHours": round(annualized_hours, 2),
        }
    return profiles


def _resolve_billing_currency(costs: list[dict[str, Any]]) -> str | None:
    currencies: list[str] = []
    for cost in costs:
        candidate = cost.get("billingCurrency") or cost.get("currency")
        if candidate:
            currencies.append(str(candidate))
    if not currencies:
        return None
    return Counter(currencies).most_common(1)[0][0]


def _lookup_vm_commitment_rates(vm_size: str, region: str) -> dict[str, Any]:
    arm_region = _normalize_region_name(region)
    sku = str(vm_size or "").strip()
    if not sku or not arm_region:
        return {
            "onDemandHourly": None,
            "reservation1YearHourly": None,
            "reservation2YearHourly": None,
            "reservation3YearHourly": None,
            "savingsPlan1YearHourly": None,
            "savingsPlan2YearHourly": None,
            "savingsPlan3YearHourly": None,
            "pricingCurrency": None,
        }

    base_filter = (
        "serviceName eq 'Virtual Machines'"
        f" and armRegionName eq '{arm_region}'"
        f" and armSkuName eq '{sku}'"
    )

    on_demand_items = _retail_prices_query(base_filter + " and priceType eq 'Consumption'")
    on_demand_hourly = _pick_min_unit_price(on_demand_items)

    reservation_1y_items = _retail_prices_query(base_filter + " and type eq 'Reservation' and reservationTerm eq '1 Year'")
    reservation_3y_items = _retail_prices_query(base_filter + " and type eq 'Reservation' and reservationTerm eq '3 Years'")
    reservation_1y_hourly = _pick_min_unit_price(reservation_1y_items)
    reservation_3y_hourly = _pick_min_unit_price(reservation_3y_items)
    reservation_2y_hourly = None
    if reservation_1y_hourly is not None and reservation_3y_hourly is not None:
        reservation_2y_hourly = round((reservation_1y_hourly + reservation_3y_hourly) / 2, 6)
    elif reservation_1y_hourly is not None:
        reservation_2y_hourly = reservation_1y_hourly
    elif reservation_3y_hourly is not None:
        reservation_2y_hourly = reservation_3y_hourly

    savings_plan_1y_hourly = None
    savings_plan_3y_hourly = None
    for item in on_demand_items:
        savings_plan = item.get("savingsPlan")
        if not isinstance(savings_plan, list):
            continue
        for term_entry in savings_plan:
            term = str(term_entry.get("term") or "")
            price = _to_float(term_entry.get("unitPrice"))
            if price is None or price <= 0:
                continue
            if term == "1 Year":
                savings_plan_1y_hourly = price if savings_plan_1y_hourly is None else min(savings_plan_1y_hourly, price)
            elif term == "3 Years":
                savings_plan_3y_hourly = price if savings_plan_3y_hourly is None else min(savings_plan_3y_hourly, price)

    savings_plan_2y_hourly = None
    if savings_plan_1y_hourly is not None and savings_plan_3y_hourly is not None:
        savings_plan_2y_hourly = round((savings_plan_1y_hourly + savings_plan_3y_hourly) / 2, 6)
    elif savings_plan_1y_hourly is not None:
        savings_plan_2y_hourly = savings_plan_1y_hourly
    elif savings_plan_3y_hourly is not None:
        savings_plan_2y_hourly = savings_plan_3y_hourly

    pricing_currency = _pick_currency_code(on_demand_items, reservation_1y_items, reservation_3y_items)

    return {
        "onDemandHourly": on_demand_hourly,
        "reservation1YearHourly": reservation_1y_hourly,
        "reservation2YearHourly": reservation_2y_hourly,
        "reservation3YearHourly": reservation_3y_hourly,
        "savingsPlan1YearHourly": savings_plan_1y_hourly,
        "savingsPlan2YearHourly": savings_plan_2y_hourly,
        "savingsPlan3YearHourly": savings_plan_3y_hourly,
        "pricingCurrency": pricing_currency,
    }


def _calc_commitment_savings_from_rates(on_demand_hourly: float | None, commitment_hourly: float | None, annualized_hours: float, years: int) -> float | None:
    if on_demand_hourly is None or commitment_hourly is None:
        return None
    savings = max(0.0, (on_demand_hourly - commitment_hourly) * annualized_hours * years)
    return round(savings, 2)


def _add_pricing_based_commitment_savings(
    recommendation: dict[str, Any],
    resource: dict[str, Any] | None,
    billing_currency: str | None,
    vm_usage_profile: dict[str, dict[str, float]],
    vm_size_by_resource_id: dict[str, str],
    rates_cache: dict[tuple[str, str], dict[str, Any]],
    fx_config: dict[str, Any],
) -> bool:
    if not resource:
        return False
    resource_type = str(resource.get("type") or "").lower()
    if resource_type != "microsoft.compute/virtualmachines":
        return False

    resource_id = str(resource.get("id") or "")
    resource_id_lower = resource_id.lower()
    vm_size = vm_size_by_resource_id.get(resource_id_lower)
    if not vm_size and resource_id:
        try:
            vm_size_result = run_az("vm", "show", "--ids", resource_id, "--query", "hardwareProfile.vmSize")
            if vm_size_result:
                vm_size = str(vm_size_result)
                vm_size_by_resource_id[resource_id_lower] = vm_size
        except Exception as exc:
            logging.warning("Unable to resolve VM size for %s: %s", resource_id, exc)
            return False

    if not vm_size:
        return False

    region = str(resource.get("location") or "")
    arm_region_name = _normalize_region_name(region)
    cache_key = (vm_size, arm_region_name)
    rates = rates_cache.get(cache_key)
    if rates is None:
        rates = _get_cached_vm_commitment_rates(vm_size, arm_region_name)
        if rates is None:
            try:
                rates = _lookup_vm_commitment_rates(vm_size, region)
                _set_cached_vm_commitment_rates(vm_size, arm_region_name, rates)
            except Exception as exc:
                logging.warning("Unable to lookup retail pricing for %s in %s: %s", vm_size, region, exc)
                return False
        rates_cache[cache_key] = rates

    on_demand_hourly = rates.get("onDemandHourly")
    usage_profile = vm_usage_profile.get(resource_id_lower, {})
    annualized_hours = float(usage_profile.get("annualizedHours") or HOURS_PER_YEAR)
    recommendation["reservationSavings1Year"] = _calc_commitment_savings_from_rates(on_demand_hourly, rates.get("reservation1YearHourly"), annualized_hours, 1)
    recommendation["reservationSavings2Year"] = _calc_commitment_savings_from_rates(on_demand_hourly, rates.get("reservation2YearHourly"), annualized_hours, 2)
    recommendation["reservationSavings3Year"] = _calc_commitment_savings_from_rates(on_demand_hourly, rates.get("reservation3YearHourly"), annualized_hours, 3)
    recommendation["savingsPlanSavings1Year"] = _calc_commitment_savings_from_rates(on_demand_hourly, rates.get("savingsPlan1YearHourly"), annualized_hours, 1)
    recommendation["savingsPlanSavings2Year"] = _calc_commitment_savings_from_rates(on_demand_hourly, rates.get("savingsPlan2YearHourly"), annualized_hours, 2)
    recommendation["savingsPlanSavings3Year"] = _calc_commitment_savings_from_rates(on_demand_hourly, rates.get("savingsPlan3YearHourly"), annualized_hours, 3)
    recommendation["savingsEstimateBasis"] = "retail-pricing-usage-scaled"
    recommendation["resourceVmSize"] = vm_size
    recommendation["pricingCurrency"] = rates.get("pricingCurrency")
    recommendation["billingCurrency"] = billing_currency
    recommendation["currencyMatch"] = (
        str(recommendation["pricingCurrency"]).upper() == str(billing_currency).upper()
        if recommendation.get("pricingCurrency") and billing_currency
        else None
    )
    recommendation["savingsCurrency"] = recommendation.get("pricingCurrency")
    recommendation["fxApplied"] = False
    recommendation["fxRateApplied"] = None
    recommendation["fxSource"] = fx_config.get("source")

    pricing_currency = recommendation.get("pricingCurrency")
    if pricing_currency and billing_currency and recommendation.get("currencyMatch") is False:
        fx_rate_applied: float | None = None
        converted_any = False
        for field in [
            "reservationSavings1Year",
            "reservationSavings2Year",
            "reservationSavings3Year",
            "savingsPlanSavings1Year",
            "savingsPlanSavings2Year",
            "savingsPlanSavings3Year",
        ]:
            converted_value, fx_rate = _convert_currency_amount(
                _to_float(recommendation.get(field)),
                pricing_currency,
                billing_currency,
                fx_config,
            )
            if converted_value is not None:
                recommendation[field] = converted_value
                converted_any = True
            if fx_rate is not None and fx_rate_applied is None:
                fx_rate_applied = fx_rate
        if converted_any:
            recommendation["savingsCurrency"] = _normalize_currency_code(billing_currency)
            recommendation["fxApplied"] = True
            recommendation["fxRateApplied"] = fx_rate_applied

    recommendation["estimatedAverageInstanceCount"] = usage_profile.get("averageInstanceCount")
    recommendation["observedUsageHoursInPeriod"] = usage_profile.get("usageHoursInPeriod")
    return True


def _normalize_subscriptions(raw_subscriptions: list[dict[str, Any]]) -> list[dict[str, Any]]:
    normalized: list[dict[str, Any]] = []
    for item in raw_subscriptions:
        if not isinstance(item, dict):
            continue
        sub_id = str(item.get("id") or item.get("subscriptionId") or "").strip()
        if not sub_id:
            continue
        managed_tenants: list[str] = []
        raw_managed = item.get("managedByTenants")
        if isinstance(raw_managed, list):
            for managed in raw_managed:
                if isinstance(managed, dict):
                    tenant_id = str(managed.get("tenantId") or "").strip()
                    if tenant_id:
                        managed_tenants.append(tenant_id)
        normalized.append(
            {
                "id": sub_id,
                "name": str(item.get("name") or item.get("displayName") or sub_id).strip(),
                "tenantId": str(item.get("tenantId") or item.get("homeTenantId") or "").strip(),
                "homeTenantId": str(item.get("homeTenantId") or "").strip(),
                "managedTenantIds": managed_tenants,
            }
        )
    return normalized


def collect_report_data(
    target_tenant_id: str | None = None,
    target_subscription_id: str | None = None,
    arm_access_token: str | None = None,
) -> dict[str, Any]:
    start_date, end_date = _month_range()
    try:
        subscriptions = run_az("account", "list", arm_access_token=arm_access_token)
    except Exception as exc:
        logging.warning("Subscription discovery via account list failed: %s", exc)
        subscriptions = []

    if not isinstance(subscriptions, list):
        raise RuntimeError("Unable to enumerate Azure subscriptions. Verify Azure CLI authentication and account access.")

    normalized_subscriptions = _normalize_subscriptions(subscriptions)

    if not normalized_subscriptions:
        fallback_sub_id = _CURRENT_SUBSCRIPTION_ID or _subscription_id_from_owner_name()
        if fallback_sub_id:
            normalized_subscriptions = [{"id": fallback_sub_id, "name": "Current subscription", "tenantId": ""}]

    all_visible_subscription_ids = {
        str(sub.get("id") or "").strip().lower() for sub in normalized_subscriptions if str(sub.get("id") or "").strip()
    }

    tenant_filter = (target_tenant_id or "").strip().lower()
    subscription_filter = (target_subscription_id or "").strip().lower()

    if tenant_filter:
        filtered_subscriptions: list[dict[str, Any]] = []
        for sub in normalized_subscriptions:
            candidate_tenants = {
                str(sub.get("tenantId") or "").strip().lower(),
                str(sub.get("homeTenantId") or "").strip().lower(),
            }
            managed_tenants = sub.get("managedTenantIds")
            if isinstance(managed_tenants, list):
                candidate_tenants.update(str(tenant_id).strip().lower() for tenant_id in managed_tenants)
            if tenant_filter in candidate_tenants:
                filtered_subscriptions.append(sub)
        normalized_subscriptions = filtered_subscriptions

    if subscription_filter:
        normalized_subscriptions = [
            sub for sub in normalized_subscriptions if str(sub.get("id") or "").strip().lower() == subscription_filter
        ]

    if not normalized_subscriptions:
        if subscription_filter and subscription_filter not in all_visible_subscription_ids:
            raise RuntimeError(
                "The requested subscription is not visible to the current identity. Run `az login` with an account that has access to that subscription, then retry."
            )
        raise RuntimeError(
            "No matching Azure subscriptions found for the requested tenant/subscription. Authenticate to that tenant and ensure the identity has access."
        )

    logging.info("Collecting FinOps report data across %d subscriptions", len(normalized_subscriptions))
    resources: list[dict[str, Any]] = []
    costs: list[dict[str, Any]] = []
    recommendations: list[dict[str, Any]] = []
    vm_size_by_resource_id: dict[str, str] = {}
    rates_cache: dict[tuple[str, str], dict[str, Any]] = {}
    fx_config = _load_fx_configuration()

    for subscription in normalized_subscriptions:
        try:
            run_az("account", "set", "--subscription", subscription["id"])
        except Exception as exc:
            logging.warning("Skipping subscription switch for %s (%s): %s", subscription.get("name"), subscription.get("id"), exc)
            continue

        try:
            current_resources = run_az(
                "resource", "list",
                "--query",
                "[].{id:id, name:name, type:type, location:location, resourceGroup:resourceGroup, subscriptionId:subscriptionId}",
                arm_access_token=arm_access_token,
            )
            if isinstance(current_resources, list):
                resources.extend(current_resources)
        except Exception as exc:
            logging.warning("Skipping resource query for %s: %s", subscription.get("name"), exc)
            continue

        try:
            current_costs = run_az(
                "consumption", "usage", "list",
                "--start-date", start_date,
                "--end-date", end_date,
                "--query",
                "[].{resourceId:instanceId, resourceName:instanceName, cost:pretaxCost, usageQuantity:usageQuantity, date:date, meterCategory:meterCategory, resourceGroup:resourceGroup, billingCurrency:billingCurrency, currency:currency}",
                arm_access_token=arm_access_token,
            )
            if isinstance(current_costs, list):
                for item in current_costs:
                    costs.append({**item, "subscriptionId": subscription["id"], "subscriptionName": subscription["name"]})
        except Exception as exc:
            logging.warning("Skipping cost query for %s: %s", subscription["name"], exc)

        try:
            current_recs = run_az(
                "advisor", "recommendation", "list",
                "--query",
                "[].{id:id, name:name, category:category, impact:impact, severity:severity, description:shortDescription.problem, solution:shortDescription.solution, resourceId:resourceMetadata.resourceId, annualSavingsAmount:annualSavingsAmount, savingsAmount:savingsAmount, extendedProperties:extendedProperties}",
                arm_access_token=arm_access_token,
            )
            if isinstance(current_recs, list):
                recommendations.extend(current_recs)
        except Exception as exc:
            logging.warning("Skipping advisor recommendations for %s: %s", subscription["name"], exc)

    billing_currency = _resolve_billing_currency(costs)
    vm_usage_profile = _build_vm_usage_profile(costs, start_date, end_date)

    resource_lookup = {resource["id"].lower(): resource for resource in resources if resource.get("id")}
    for recommendation in recommendations:
        resource = resource_lookup.get(str(recommendation.get("resourceId", "")).lower(), {})
        fallback_name, fallback_type = _resource_fallback_from_id(recommendation.get("resourceId"))
        recommendation["resourceName"] = resource.get("name") or fallback_name
        recommendation["resourceType"] = resource.get("type") or fallback_type
        if not _add_pricing_based_commitment_savings(
            recommendation,
            resource,
            billing_currency,
            vm_usage_profile,
            vm_size_by_resource_id,
            rates_cache,
            fx_config,
        ):
            _add_commitment_savings_fields(recommendation)

    return {
        "generatedAt": datetime.now(timezone.utc).isoformat(),
        "startDate": start_date,
        "endDate": end_date,
        "tenantId": target_tenant_id,
        "subscriptionId": target_subscription_id,
        "resources": resources,
        "costs": costs,
        "recommendations": recommendations,
    }


def write_sqlite(snapshot: dict[str, Any]) -> Path:
    db_path = REPORT_DB_PATH
    conn = sqlite3.connect(db_path)
    _ensure_retail_cache_table(conn)
    conn.execute("DROP TABLE IF EXISTS resource_costs")
    conn.execute("DROP TABLE IF EXISTS recommendations")
    conn.execute("DROP TABLE IF EXISTS resources")
    conn.execute("DROP TABLE IF EXISTS report_runs")

    conn.execute("CREATE TABLE report_runs (generatedAt TEXT, startDate TEXT, endDate TEXT)")
    conn.execute("CREATE TABLE resources (id TEXT, name TEXT, type TEXT, location TEXT, resourceGroup TEXT, subscriptionId TEXT)")
    conn.execute("CREATE TABLE resource_costs (resourceId TEXT, resourceName TEXT, subscriptionName TEXT, cost REAL, usageQuantity REAL, billingCurrency TEXT, currency TEXT, date TEXT, meterCategory TEXT)")
    conn.execute("CREATE TABLE recommendations (id TEXT, name TEXT, category TEXT, impact TEXT, severity TEXT, description TEXT, solution TEXT, resourceId TEXT, resourceName TEXT, resourceType TEXT, resourceVmSize TEXT, savingsEstimateBasis TEXT, pricingCurrency TEXT, billingCurrency TEXT, currencyMatch TEXT, savingsCurrency TEXT, fxApplied TEXT, fxRateApplied REAL, fxSource TEXT, estimatedAverageInstanceCount REAL, observedUsageHoursInPeriod REAL, reservationSavings1Year REAL, reservationSavings2Year REAL, reservationSavings3Year REAL, savingsPlanSavings1Year REAL, savingsPlanSavings2Year REAL, savingsPlanSavings3Year REAL)")

    conn.execute("INSERT INTO report_runs VALUES (?, ?, ?)", (snapshot["generatedAt"], snapshot["startDate"], snapshot["endDate"]))
    for resource in snapshot["resources"]:
        conn.execute("INSERT INTO resources VALUES (?, ?, ?, ?, ?, ?)", (resource.get("id"), resource.get("name"), resource.get("type"), resource.get("location"), resource.get("resourceGroup"), resource.get("subscriptionId")))
    for cost in snapshot["costs"]:
        conn.execute("INSERT INTO resource_costs VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)", (cost.get("resourceId"), cost.get("resourceName"), cost.get("subscriptionName"), float(cost.get("cost") or 0), _to_float(cost.get("usageQuantity")), cost.get("billingCurrency"), cost.get("currency"), cost.get("date"), cost.get("meterCategory")))
    for rec in snapshot["recommendations"]:
        conn.execute(
            "INSERT INTO recommendations VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)",
            (
                rec.get("id"),
                rec.get("name"),
                rec.get("category"),
                rec.get("impact"),
                rec.get("severity"),
                rec.get("description"),
                rec.get("solution"),
                rec.get("resourceId"),
                rec.get("resourceName"),
                rec.get("resourceType"),
                rec.get("resourceVmSize"),
                rec.get("savingsEstimateBasis"),
                rec.get("pricingCurrency"),
                rec.get("billingCurrency"),
                None if rec.get("currencyMatch") is None else str(bool(rec.get("currencyMatch"))),
                rec.get("savingsCurrency"),
                None if rec.get("fxApplied") is None else str(bool(rec.get("fxApplied"))),
                _to_float(rec.get("fxRateApplied")),
                rec.get("fxSource"),
                _to_float(rec.get("estimatedAverageInstanceCount")),
                _to_float(rec.get("observedUsageHoursInPeriod")),
                _to_float(rec.get("reservationSavings1Year")),
                _to_float(rec.get("reservationSavings2Year")),
                _to_float(rec.get("reservationSavings3Year")),
                _to_float(rec.get("savingsPlanSavings1Year")),
                _to_float(rec.get("savingsPlanSavings2Year")),
                _to_float(rec.get("savingsPlanSavings3Year")),
            ),
        )
    conn.commit()
    conn.close()
    return db_path


def make_pdf(snapshot: dict[str, Any]) -> Path:
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import landscape, letter
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.units import inch
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
    except Exception as exc:
        raise RuntimeError("PDF generation dependencies are unavailable") from exc

    pdf_path = ARTIFACT_DIR / "finops_summary.pdf"
    styles = getSampleStyleSheet()
    body_style = styles["BodyText"].clone("FinOpsBody")
    body_style.fontSize = 8
    body_style.leading = 9
    heading1_style = styles["Heading1"].clone("FinOpsHeading1")
    heading1_style.fontSize = 15
    heading1_style.leading = 17
    heading2_style = styles["Heading2"].clone("FinOpsHeading2")
    heading2_style.fontSize = 10
    heading2_style.leading = 12
    cell_style = styles["BodyText"].clone("FinOpsCell")
    cell_style.fontSize = 5.5
    cell_style.leading = 6.2
    cell_style.wordWrap = "CJK"
    cell_style.splitLongWords = True
    header_cell_style = styles["BodyText"].clone("FinOpsHeaderCell")
    header_cell_style.fontSize = 6
    header_cell_style.leading = 6.6
    header_cell_style.wordWrap = "CJK"
    header_cell_style.splitLongWords = True

    def _pdf_cell(value: Any, *, header: bool = False) -> Paragraph:
        text = "" if value is None else str(value)
        text = escape(text).replace("\n", "<br/>")
        return Paragraph(text, header_cell_style if header else cell_style)

    story = [Paragraph("Azure FinOps Summary", heading1_style), Spacer(1, 0.12 * inch)]
    story.append(Paragraph(f"Generated: {snapshot['generatedAt']}", body_style))
    story.append(Paragraph(f"Period: {snapshot['startDate']} to {snapshot['endDate']}", body_style))
    story.append(Spacer(1, 0.2 * inch))

    total_cost = round(sum(float(item.get('cost') or 0) for item in snapshot['costs']), 2)
    story.append(Paragraph(f"Total cost in period: ${total_cost:,.2f}", heading2_style))
    story.append(Spacer(1, 0.1 * inch))

    resource_table = [[
        _pdf_cell('Resource', header=True),
        _pdf_cell('Type', header=True),
        _pdf_cell('Location', header=True),
        _pdf_cell('Subscription', header=True),
    ], *[
        (
            _pdf_cell(r.get('name', '-')),
            _pdf_cell(r.get('type', '-')),
            _pdf_cell(r.get('location', '-')),
            _pdf_cell(r.get('subscriptionId', '-')),
        )
        for r in snapshot['resources'][:15]
    ]]
    table = Table(resource_table, repeatRows=1, colWidths=[2.5 * inch, 2.3 * inch, 1.6 * inch, 2.4 * inch])
    table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2F6FED')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
        ('FONTNAME', (0, 0), (-1, -1), 'Helvetica'),
        ('FONTSIZE', (0, 0), (-1, -1), 7),
        ('LEADING', (0, 0), (-1, -1), 8),
        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
        ('LEFTPADDING', (0, 0), (-1, -1), 3),
        ('RIGHTPADDING', (0, 0), (-1, -1), 3),
        ('TOPPADDING', (0, 0), (-1, -1), 2),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 2),
    ]))
    story.append(Paragraph("Top resources", heading2_style))
    story.append(table)
    story.append(Spacer(1, 0.2 * inch))

    reservation_1y_total = round(sum(_to_float(rec.get('reservationSavings1Year')) or 0 for rec in snapshot['recommendations']), 2)
    reservation_3y_total = round(sum(_to_float(rec.get('reservationSavings3Year')) or 0 for rec in snapshot['recommendations']), 2)
    savings_plan_1y_total = round(sum(_to_float(rec.get('savingsPlanSavings1Year')) or 0 for rec in snapshot['recommendations']), 2)
    savings_plan_3y_total = round(sum(_to_float(rec.get('savingsPlanSavings3Year')) or 0 for rec in snapshot['recommendations']), 2)

    story.append(Spacer(1, 0.2 * inch))
    story.append(Paragraph("Estimated commitment savings totals", heading2_style))
    story.append(Paragraph(f"Reservation 1Y: ${reservation_1y_total:,.2f}", body_style))
    story.append(Paragraph(f"Reservation 3Y: ${reservation_3y_total:,.2f}", body_style))
    story.append(Paragraph(f"Savings Plan 1Y: ${savings_plan_1y_total:,.2f}", body_style))
    story.append(Paragraph(f"Savings Plan 3Y: ${savings_plan_3y_total:,.2f}", body_style))
    story.append(Spacer(1, 0.1 * inch))

    rec_table = [[
        _pdf_cell('Recommendation', header=True),
        _pdf_cell('Impact', header=True),
        _pdf_cell('Resource', header=True),
        _pdf_cell('Resv 1Y', header=True),
        _pdf_cell('Resv 2Y', header=True),
        _pdf_cell('Resv 3Y', header=True),
        _pdf_cell('SP 1Y', header=True),
        _pdf_cell('SP 2Y', header=True),
        _pdf_cell('SP 3Y', header=True),
        _pdf_cell('Currency', header=True),
        _pdf_cell('Basis', header=True),
    ], *[
        (
            _pdf_cell(rec.get('name', '-')),
            _pdf_cell(rec.get('impact', '-')),
            _pdf_cell(rec.get('resourceName', '-')),
            _pdf_cell(rec.get('reservationSavings1Year', '-')),
            _pdf_cell(rec.get('reservationSavings2Year', '-')),
            _pdf_cell(rec.get('reservationSavings3Year', '-')),
            _pdf_cell(rec.get('savingsPlanSavings1Year', '-')),
            _pdf_cell(rec.get('savingsPlanSavings2Year', '-')),
            _pdf_cell(rec.get('savingsPlanSavings3Year', '-')),
            _pdf_cell(rec.get('savingsCurrency') or rec.get('pricingCurrency') or '-'),
            _pdf_cell(rec.get('savingsEstimateBasis', '-')),
        )
        for rec in snapshot['recommendations']
    ]]
    rec = Table(
        rec_table,
        repeatRows=1,
        colWidths=[1.55 * inch, 0.5 * inch, 1.15 * inch, 0.55 * inch, 0.55 * inch, 0.55 * inch, 0.55 * inch, 0.55 * inch, 0.55 * inch, 0.55 * inch, 0.65 * inch],
    )
    rec.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#0F766E')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
        ('FONTNAME', (0, 0), (-1, -1), 'Helvetica'),
        ('FONTSIZE', (0, 0), (-1, -1), 6),
        ('LEADING', (0, 0), (-1, -1), 7),
        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
        ('LEFTPADDING', (0, 0), (-1, -1), 2),
        ('RIGHTPADDING', (0, 0), (-1, -1), 2),
        ('TOPPADDING', (0, 0), (-1, -1), 1),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 1),
    ]))
    story.append(Paragraph("Advisor recommendations", heading2_style))
    story.append(rec)

    doc = SimpleDocTemplate(
        str(pdf_path),
        pagesize=landscape(letter),
        leftMargin=0.3 * inch,
        rightMargin=0.3 * inch,
        topMargin=0.35 * inch,
        bottomMargin=0.35 * inch,
    )
    doc.build(story)
    return pdf_path


def make_powerpoint(snapshot: dict[str, Any]) -> Path:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except Exception as exc:
        raise RuntimeError("PowerPoint generation dependencies are unavailable") from exc

    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)

    recommendations = snapshot.get('recommendations', [])
    resources = snapshot.get('resources', [])

    resource_lookup = {
        str(resource.get('id') or '').lower(): resource
        for resource in resources
        if resource.get('id')
    }
    associated_resources: dict[str, dict[str, Any]] = {}
    category_counts: Counter[str] = Counter()
    impact_counts: Counter[str] = Counter()

    for rec in recommendations:
        category = str(rec.get('category') or 'Unknown')
        impact = str(rec.get('impact') or 'Unknown')
        category_counts[category] += 1
        impact_counts[impact] += 1

        resource_id = str(rec.get('resourceId') or '').strip()
        if not resource_id:
            continue

        key = resource_id.lower()
        if key not in associated_resources:
            resource_info = resource_lookup.get(key, {})
            associated_resources[key] = {
                'resourceId': resource_id,
                'name': resource_info.get('name') or rec.get('resourceName') or '-',
                'type': resource_info.get('type') or rec.get('resourceType') or '-',
                'location': resource_info.get('location') or '-',
                'resourceGroup': resource_info.get('resourceGroup') or '-',
                'subscriptionId': resource_info.get('subscriptionId') or '-',
                'recommendationCount': 0,
                'highImpactCount': 0,
                'reservationSavings1YearTotal': 0.0,
                'savingsPlanSavings1YearTotal': 0.0,
            }

        current = associated_resources[key]
        current['recommendationCount'] += 1
        if impact.lower() == 'high':
            current['highImpactCount'] += 1
        current['reservationSavings1YearTotal'] += _to_float(rec.get('reservationSavings1Year')) or 0.0
        current['savingsPlanSavings1YearTotal'] += _to_float(rec.get('savingsPlanSavings1Year')) or 0.0

    associated_resource_rows = sorted(
        associated_resources.values(),
        key=lambda item: (
            -int(item.get('recommendationCount') or 0),
            str(item.get('name') or '').lower(),
        ),
    )

    reservation_1y_total = round(sum(_to_float(rec.get('reservationSavings1Year')) or 0 for rec in recommendations), 2)
    reservation_3y_total = round(sum(_to_float(rec.get('reservationSavings3Year')) or 0 for rec in recommendations), 2)
    savings_plan_1y_total = round(sum(_to_float(rec.get('savingsPlanSavings1Year')) or 0 for rec in recommendations), 2)
    savings_plan_3y_total = round(sum(_to_float(rec.get('savingsPlanSavings3Year')) or 0 for rec in recommendations), 2)

    title_color = RGBColor(17, 24, 39)
    subtitle_color = RGBColor(71, 85, 105)
    body_color = RGBColor(30, 41, 59)
    header_fill_color = RGBColor(30, 64, 175)
    header_text_color = RGBColor(255, 255, 255)
    row_alt_fill_color = RGBColor(248, 250, 252)
    background_color = RGBColor(245, 247, 250)

    def _set_slide_background(slide_obj: Any) -> None:
        fill = slide_obj.background.fill
        fill.solid()
        fill.fore_color.rgb = background_color

    def _style_title(slide_obj: Any, size: int = 28) -> None:
        title_shape = slide_obj.shapes.title
        if not title_shape or not title_shape.has_text_frame:
            return
        title_frame = title_shape.text_frame
        for paragraph in title_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.name = "Segoe UI Semibold"
                run.font.size = Pt(size)
                run.font.color.rgb = title_color

    def _set_text_frame_lines(text_frame: Any, lines: list[str], font_size: int) -> None:
        text_frame.clear()
        for idx, line in enumerate(lines):
            paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
            paragraph.text = line
            paragraph.level = 0
            paragraph.space_after = Pt(6)
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.name = "Segoe UI"
                run.font.size = Pt(font_size)
                run.font.color.rgb = subtitle_color if idx == 0 else body_color

    slide = deck.slides.add_slide(deck.slide_layouts[1])
    _set_slide_background(slide)
    slide.shapes.title.text = "Azure FinOps Summary"
    _style_title(slide, size=30)
    _set_text_frame_lines(
        slide.placeholders[1].text_frame,
        [
            f"Generated {snapshot['generatedAt']}",
            f"Period {snapshot['startDate']} to {snapshot['endDate']}",
            f"Advisor recommendations: {len(recommendations)}",
            f"Resources with recommendations: {len(associated_resource_rows)}",
        ],
        font_size=15,
    )

    slide2 = deck.slides.add_slide(deck.slide_layouts[1])
    _set_slide_background(slide2)
    slide2.shapes.title.text = "Advisor recommendations summary"
    _style_title(slide2, size=26)
    body = slide2.placeholders[1].text_frame
    top_categories = category_counts.most_common(3)
    top_impacts = impact_counts.most_common(3)
    summary_lines = [
        f"Top categories: {', '.join(f'{name} ({count})' for name, count in top_categories) if top_categories else '-'}",
        f"By impact: {', '.join(f'{name} ({count})' for name, count in top_impacts) if top_impacts else '-'}",
        f"Reservation savings 1Y: {reservation_1y_total:,.2f}",
        f"Reservation savings 3Y: {reservation_3y_total:,.2f}",
        f"Savings Plan savings 1Y: {savings_plan_1y_total:,.2f}",
        f"Savings Plan savings 3Y: {savings_plan_3y_total:,.2f}",
    ]
    _set_text_frame_lines(body, summary_lines, font_size=14)

    slide3 = deck.slides.add_slide(deck.slide_layouts[1])
    _set_slide_background(slide3)
    slide3.shapes.title.text = "Top Advisor opportunities"
    _style_title(slide3, size=26)
    body3 = slide3.placeholders[1].text_frame
    top_recommendations = sorted(
        recommendations,
        key=lambda rec: (
            _impact_rank(rec.get('impact')),
            -(_to_float(rec.get('reservationSavings1Year')) or 0),
            -(_to_float(rec.get('savingsPlanSavings1Year')) or 0),
        ),
    )
    top_lines: list[str] = []
    for rec in top_recommendations[:6]:
        top_lines.append(
            f"{rec.get('name', '-')} ({rec.get('impact', '-')}) | {rec.get('resourceName', '-')} | "
            f"Resv1Y {rec.get('reservationSavings1Year', '-')} | "
            f"SP1Y {rec.get('savingsPlanSavings1Year', '-')}"
        )
    _set_text_frame_lines(body3, top_lines or ["No advisor recommendations found in this period."], font_size=13)

    slide4 = deck.slides.add_slide(deck.slide_layouts[5])
    _set_slide_background(slide4)
    slide4.shapes.title.text = "Resources associated with Advisor recommendations"
    _style_title(slide4, size=24)
    top = Inches(1.2)
    left = Inches(0.3)
    width = Inches(12.7)
    height = Inches(5.8)
    rows = min(len(associated_resource_rows), 12) + 1
    cols = 8
    table_shape = slide4.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    headers = ["Resource Name", "Resource Type", "Location", "Resource Group", "Recommendations", "High Impact", "Resv1Y Total", "SP1Y Total"]
    column_widths = [
        Inches(2.2),
        Inches(2.4),
        Inches(1.0),
        Inches(1.4),
        Inches(1.0),
        Inches(0.9),
        Inches(1.4),
        Inches(1.4),
    ]
    for idx, column_width in enumerate(column_widths):
        table.columns[idx].width = column_width

    def _format_ppt_cell(cell: Any, text: str, is_header: bool = False) -> None:
        cell.text = text
        cell.margin_left = Inches(0.03)
        cell.margin_right = Inches(0.03)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = header_fill_color if is_header else RGBColor(255, 255, 255)
        text_frame = cell.text_frame
        text_frame.word_wrap = True
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.name = "Segoe UI"
                run.font.size = Pt(8 if is_header else 7)
                run.font.color.rgb = header_text_color if is_header else body_color

    for idx, header in enumerate(headers):
        _format_ppt_cell(table.cell(0, idx), header, is_header=True)

    for row_idx, resource in enumerate(associated_resource_rows[: rows - 1], start=1):
        values = [
            str(resource.get('name', '-')),
            str(resource.get('type', '-')),
            str(resource.get('location', '-')),
            str(resource.get('resourceGroup', '-')),
            str(resource.get('recommendationCount', 0)),
            str(resource.get('highImpactCount', 0)),
            f"{float(resource.get('reservationSavings1YearTotal') or 0):,.2f}",
            f"{float(resource.get('savingsPlanSavings1YearTotal') or 0):,.2f}",
        ]
        for col_idx, value in enumerate(values):
            _format_ppt_cell(table.cell(row_idx, col_idx), value)
            if row_idx % 2 == 0:
                alt_fill = table.cell(row_idx, col_idx).fill
                alt_fill.solid()
                alt_fill.fore_color.rgb = row_alt_fill_color

    pptx_path = ARTIFACT_DIR / "finops_summary.pptx"
    deck.save(pptx_path)
    return pptx_path


def generate_artifacts(snapshot: dict[str, Any]) -> dict[str, str]:
    _cleanup_old_artifacts()
    db_path = write_sqlite(snapshot)
    pdf_path = make_pdf(snapshot)
    pptx_path = make_powerpoint(snapshot)
    return {
        "database": str(db_path),
        "pdf": str(pdf_path),
        "powerpoint": str(pptx_path),
    }


def _build_artifact_urls(req: func.HttpRequest) -> dict[str, str]:
    base_url = req.url.split("/api/", 1)[0] + "/api"
    code = req.params.get("code")
    suffix = f"&code={quote(code)}" if code else ""
    return {
        "database": f"{base_url}/finops_artifact?name=database{suffix}",
        "pdf": f"{base_url}/finops_artifact?name=pdf{suffix}",
        "powerpoint": f"{base_url}/finops_artifact?name=powerpoint{suffix}",
    }


def _to_rows(items: list[dict[str, Any]], columns: list[str]) -> list[list[Any]]:
    rows: list[list[Any]] = []
    for item in items:
        rows.append([item.get(column) for column in columns])
    return rows


def _parse_positive_int(value: str | None, default: int) -> int:
    if value is None:
        return default
    try:
        parsed = int(value)
        return parsed if parsed > 0 else default
    except (TypeError, ValueError):
        return default


def _impact_rank(impact: Any) -> int:
    mapping = {
        "High": 0,
        "Medium": 1,
        "Low": 2,
    }
    return mapping.get(str(impact), 3)


def _looks_like_opaque_identifier(value: Any) -> bool:
    text = str(value or "").strip()
    if not text:
        return False

    # Typical GUID shape.
    if len(text) == 36 and text.count("-") == 4:
        parts = text.split("-")
        part_lengths = [8, 4, 4, 4, 12]
        if len(parts) == 5 and all(len(part) == expected for part, expected in zip(parts, part_lengths)):
            return True

    # Advisor often emits opaque hash-like recommendation names.
    if len(text) in (32, 40, 64):
        lowered = text.lower()
        if all(ch in "0123456789abcdef" for ch in lowered):
            return True

    return False


def _recommendation_display_name(rec: dict[str, Any]) -> str:
    resource_name = str(rec.get("resourceName") or "").strip()
    rec_name = str(rec.get("name") or "").strip()
    if _looks_like_opaque_identifier(rec_name):
        if resource_name and resource_name.lower() != "unknown resource":
            return resource_name

        resource_id = str(rec.get("resourceId") or "").strip()
        if resource_id:
            segments = [segment for segment in resource_id.split("/") if segment]
            if len(segments) >= 8:
                return segments[-1]
            if len(segments) == 2 and segments[0].lower() == "subscriptions":
                return f"Subscription {segments[1][:8]}"

        description = str(rec.get("description") or "").strip()
        if description:
            return description

    return rec_name or resource_name or "-"


def _resource_fallback_from_id(resource_id: Any) -> tuple[str, str]:
    rid = str(resource_id or "").strip()
    if not rid:
        return ("Not specified", "Unknown")

    segments = [segment for segment in rid.split("/") if segment]
    lowered = [segment.lower() for segment in segments]

    if len(segments) == 2 and lowered[0] == "subscriptions":
        sub_id = segments[1]
        return (f"Subscription {sub_id[:8]}", "Microsoft.Resources/subscriptions")

    if len(segments) >= 8:
        provider_namespace = segments[5]
        resource_type = segments[6]
        resource_name = segments[-1]
        return (resource_name, f"{provider_namespace}/{resource_type}")

    leaf = segments[-1] if segments else rid
    return (leaf, "Unknown")


def _format_usd(value: Any) -> str | None:
    amount = _to_float(value)
    if amount is None:
        return None
    return f"${amount:,.2f}"


def _resource_type_parts(resource_type: Any) -> tuple[str, str]:
    value = str(resource_type or "").strip()
    if not value or "/" not in value:
        return ("Unknown", "Unknown")
    category, sub_category = value.split("/", 1)
    return (category or "Unknown", sub_category or "Unknown")


def _paginate(items: list[dict[str, Any]], page: int, page_size: int) -> tuple[list[dict[str, Any]], int, int]:
    total_rows = len(items)
    start = (page - 1) * page_size
    end = start + page_size
    if start >= total_rows:
        return [], total_rows, start
    return items[start:end], total_rows, start


def build_tabular_output(
    snapshot: dict[str, Any],
    artifacts: dict[str, str],
    artifact_urls: dict[str, str],
    page: int,
    page_size: int,
) -> dict[str, Any]:
    resource_columns = ["name", "type", "location", "resourceGroup", "subscriptionId"]
    recommendation_columns = [
        "name",
        "description",
        "reservationSavings1Year",
        "reservationSavings2Year",
        "reservationSavings3Year",
        "savingsPlanSavings1Year",
        "savingsPlanSavings2Year",
        "savingsPlanSavings3Year",
        "savingsCurrency",
        "category",
        "impact",
        "solution",
        "resourceCategory",
        "resourceSubCategory",
        "resourceType",
        "resourceId",
        "resourceVmSize",
        "savingsEstimateBasis",
        "pricingCurrency",
        "billingCurrency",
        "currencyMatch",
        "fxApplied",
        "fxRateApplied",
        "fxSource",
        "estimatedAverageInstanceCount",
        "observedUsageHoursInPeriod",
        "recommendationName",
    ]

    resources = sorted(
        snapshot.get("resources", []),
        key=lambda r: (
            str(r.get("name") or "").lower(),
            str(r.get("type") or "").lower(),
        ),
    )
    recommendations_raw = snapshot.get("recommendations", [])
    recommendations = sorted(
        recommendations_raw,
        key=lambda r: (
            _impact_rank(r.get("impact")),
            str(r.get("category") or "").lower(),
            str(r.get("name") or "").lower(),
        ),
    )
    recommendations_for_table = []
    money_columns = {
        "reservationSavings1Year",
        "reservationSavings2Year",
        "reservationSavings3Year",
        "savingsPlanSavings1Year",
        "savingsPlanSavings2Year",
        "savingsPlanSavings3Year",
    }
    for rec in recommendations:
        row = dict(rec)
        row["recommendationName"] = rec.get("name")
        row["name"] = _recommendation_display_name(rec)
        row["savingsCurrency"] = "USD"
        resource_category, resource_sub_category = _resource_type_parts(row.get("resourceType"))
        row["resourceCategory"] = resource_category
        row["resourceSubCategory"] = resource_sub_category
        for column in money_columns:
            row[column] = _format_usd(row.get(column))
        recommendations_for_table.append(row)

    paged_resources, total_resource_rows, resource_offset = _paginate(resources, page, page_size)
    paged_recommendations, total_recommendation_rows, recommendation_offset = _paginate(recommendations_for_table, page, page_size)

    return {
        "status": "ok",
        "generatedAt": snapshot.get("generatedAt"),
        "period": {
            "startDate": snapshot.get("startDate"),
            "endDate": snapshot.get("endDate"),
        },
        "scope": {
            "tenantId": snapshot.get("tenantId") or "auto",
            "subscriptionId": snapshot.get("subscriptionId") or "auto",
        },
        "artifacts": artifacts,
        "artifactUrls": artifact_urls,
        "tabs": [
            {
                "name": "Resource Information",
                "layout": {
                    "type": "table",
                    "showGridLines": True,
                    "showColumnHeaders": True,
                },
                "columns": resource_columns,
                "rows": _to_rows(paged_resources, resource_columns),
                "page": page,
                "pageSize": page_size,
                "totalRows": total_resource_rows,
                "rowCount": len(paged_resources),
                "offset": resource_offset,
            },
            {
                "name": "Advisor Recommendations",
                "layout": {
                    "type": "table",
                    "showGridLines": True,
                    "showColumnHeaders": True,
                },
                "columns": recommendation_columns,
                "rows": _to_rows(paged_recommendations, recommendation_columns),
                "page": page,
                "pageSize": page_size,
                "totalRows": total_recommendation_rows,
                "rowCount": len(paged_recommendations),
                "offset": recommendation_offset,
            },
        ],
    }


def render_tabular_html(response: dict[str, Any]) -> str:
    def render_table(tab: dict[str, Any], tab_id: str) -> str:
        headers = "".join(f"<th>{escape(str(col))}</th>" for col in tab.get("columns", []))
        body_rows = []
        for row in tab.get("rows", []):
            cells = "".join(f"<td>{escape('' if value is None else str(value))}</td>" for value in row)
            body_rows.append(f"<tr>{cells}</tr>")
        body = "".join(body_rows) if body_rows else "<tr><td colspan='100%'>No rows</td></tr>"
        pager = (
            f"<div class='meta'>"
            f"Page {tab.get('page')} | Page Size {tab.get('pageSize')} | "
            f"Rows {tab.get('rowCount')} of {tab.get('totalRows')}"
            f"</div>"
        )
        return (
            f"<section class='tab' id='{tab_id}'>"
            f"<h2>{escape(str(tab.get('name', 'Table')))}</h2>"
            f"{pager}"
            f"<div class='top-scroll' aria-hidden='true'><div class='top-scroll-inner'></div></div>"
            f"<div class='table-wrap'><table><thead><tr>{headers}</tr></thead><tbody>{body}</tbody></table></div>"
            f"</section>"
        )

    tabs = response.get("tabs", [])
    nav_buttons = []
    sections = []
    for idx, tab in enumerate(tabs):
        tab_id = f"tab-{idx}"
        active = "active" if idx == 0 else ""
        nav_buttons.append(f"<button class='tab-btn {active}' data-target='{tab_id}'>{escape(str(tab.get('name', f'Tab {idx + 1}')))}</button>")
        section_html = render_table(tab, tab_id)
        if idx != 0:
            section_html = section_html.replace("class='tab'", "class='tab hidden'", 1)
        sections.append(section_html)

    return f"""
<!doctype html>
<html lang='en'>
<head>
    <meta charset='utf-8' />
    <meta name='viewport' content='width=device-width, initial-scale=1' />
    <title>FinOps Report</title>
    <script src='https://alcdn.msauth.net/browser/2.39.0/js/msal-browser.min.js'></script>
    <style>
        :root {{
            --bg: #eef2ff;
            --bg-alt: #dbeafe;
            --panel: #ffffff;
            --panel-strong: #f8fafc;
            --line: #cbd5e1;
            --text: #0f172a;
            --muted: #475569;
            --accent: #0f766e;
            --accent-strong: #0b5f5a;
            --accent-soft: #ccfbf1;
            --head-bg: #082f49;
            --head-text: #e0f2fe;
            --shadow: 0 14px 34px rgba(15, 23, 42, 0.12);
        }}
        body {{
            margin: 0;
            font-family: "Aptos", "Segoe UI", "Trebuchet MS", sans-serif;
            color: var(--text);
            background:
                radial-gradient(circle at 10% -20%, var(--bg-alt), transparent 45%),
                radial-gradient(circle at 90% 0%, #fef9c3, transparent 30%),
                linear-gradient(180deg, var(--bg) 0%, #f8fafc 100%);
        }}
        .container {{ max-width: 1400px; margin: 24px auto; padding: 0 16px; }}
        h1 {{ margin: 0 0 8px; font-size: 30px; letter-spacing: 0.2px; }}
        a {{ text-decoration: none; color: #155e75; }}
        .sub {{ color: var(--muted); margin-bottom: 16px; }}
        .artifacts {{ display: flex; gap: 10px; flex-wrap: wrap; margin-bottom: 14px; }}
        .artifact-link {{
            border: 1px solid #99f6e4;
            background: linear-gradient(180deg, #ffffff 0%, var(--accent-soft) 100%);
            padding: 8px 12px;
            border-radius: 999px;
            font-weight: 600;
            transition: transform 0.15s ease, box-shadow 0.15s ease;
        }}
        .artifact-link:hover {{ transform: translateY(-1px); box-shadow: 0 6px 16px rgba(15, 118, 110, 0.2); }}
        .tabs-nav {{ display: flex; gap: 8px; margin: 12px 0 16px; flex-wrap: wrap; }}
        .scope-actions {{ display: flex; align-items: center; gap: 10px; margin: 6px 0 14px; flex-wrap: wrap; }}
        .scope-btn {{
            border: 1px solid var(--line);
            background: #ffffff;
            color: var(--text);
            border-radius: 10px;
            padding: 8px 12px;
            font-weight: 600;
            cursor: pointer;
        }}
        .scope-btn:hover {{ background: #f1f5f9; border-color: #94a3b8; }}
        .scope-help {{ color: var(--muted); font-size: 12px; }}
        .scope-login {{
            border: 1px solid #7dd3fc;
            background: linear-gradient(180deg, #f0f9ff 0%, #e0f2fe 100%);
            color: #0c4a6e;
            border-radius: 10px;
            padding: 8px 12px;
            font-weight: 700;
            cursor: pointer;
        }}
        .scope-login:hover {{ background: linear-gradient(180deg, #e0f2fe 0%, #bae6fd 100%); }}
        .scope-logout {{
            border: 1px solid #cbd5e1;
            background: #ffffff;
            color: #334155;
            border-radius: 10px;
            padding: 8px 12px;
            font-weight: 600;
            cursor: pointer;
        }}
        .scope-status {{
            color: #0f172a;
            background: #f8fafc;
            border: 1px solid #e2e8f0;
            border-radius: 999px;
            padding: 4px 10px;
            font-size: 12px;
            max-width: min(60ch, 100%);
            overflow: hidden;
            text-overflow: ellipsis;
            white-space: nowrap;
        }}
        .tab-btn {{
            border: 1px solid var(--line);
            background: var(--panel);
            color: var(--text);
            padding: 8px 14px;
            border-radius: 10px;
            cursor: pointer;
            font-weight: 600;
            transition: all 0.16s ease;
        }}
        .tab-btn:hover {{ border-color: #94a3b8; background: #f8fafc; }}
        .tab-btn.active {{
            background: linear-gradient(145deg, var(--accent) 0%, var(--accent-strong) 100%);
            border-color: var(--accent-strong);
            color: #ecfeff;
            box-shadow: 0 6px 14px rgba(15, 118, 110, 0.26);
        }}
        .tab {{
            background: linear-gradient(180deg, var(--panel) 0%, var(--panel-strong) 100%);
            border: 1px solid var(--line);
            border-radius: 14px;
            padding: 14px;
            box-shadow: var(--shadow);
        }}
        .hidden {{ display: none; }}
        .meta {{ color: var(--muted); font-size: 13px; margin-bottom: 10px; }}
        .top-scroll {{ overflow-x: auto; overflow-y: hidden; border: 1px solid var(--line); border-radius: 8px; background: var(--panel); margin-bottom: 8px; }}
        .top-scroll-inner {{ height: 1px; }}
        .table-wrap {{ overflow: auto; border: 1px solid var(--line); border-radius: 10px; background: #ffffff; }}
        table {{ width: max-content; min-width: 100%; border-collapse: collapse; font-size: 13px; table-layout: auto; }}
        tr th, tr td {{ border: 1px solid #e2e8f0; padding: 8px 10px; text-align: left; vertical-align: top; min-width: clamp(120px, 14vw, 260px); white-space: normal; overflow-wrap: anywhere; }}
        tr th {{
            background: linear-gradient(180deg, var(--head-bg) 0%, #0c4a6e 100%);
            color: var(--head-text);
            position: sticky;
            top: 0;
            z-index: 1;
            font-weight: 700;
        }}
        tbody tr:nth-child(even) {{ background: #f8fafc; }}
        tbody tr:hover {{ background: #ecfeff; }}
        dialog.scope-dialog {{
            border: 1px solid var(--line);
            border-radius: 12px;
            box-shadow: 0 24px 50px rgba(15, 23, 42, 0.28);
            width: min(560px, 94vw);
            padding: 16px;
            font-family: "Aptos", "Segoe UI", "Trebuchet MS", sans-serif;
        }}
        dialog.scope-dialog::backdrop {{ background: rgba(2, 6, 23, 0.45); }}
        .scope-grid {{ display: grid; gap: 10px; margin: 10px 0 12px; }}
        .scope-grid label {{ font-size: 12px; color: var(--muted); font-weight: 600; }}
        .scope-grid input {{
            width: 100%;
            box-sizing: border-box;
            margin-top: 4px;
            border: 1px solid var(--line);
            border-radius: 8px;
            padding: 8px 10px;
            font-size: 13px;
        }}
        .scope-dialog-actions {{ display: flex; justify-content: flex-end; gap: 8px; }}
        .scope-dialog-actions button {{
            border: 1px solid var(--line);
            background: #ffffff;
            border-radius: 8px;
            padding: 7px 11px;
            cursor: pointer;
            font-weight: 600;
        }}
        .scope-dialog-actions .apply {{
            background: var(--accent);
            border-color: var(--accent-strong);
            color: #ecfeff;
        }}
        @media (max-width: 900px) {{
            table {{ font-size: 12px; }}
            tr th, tr td {{ padding: 7px 8px; min-width: clamp(100px, 24vw, 220px); }}
        }}
        @media (max-width: 600px) {{
            h1 {{ font-size: 22px; }}
            .container {{ margin: 16px auto; padding: 0 10px; }}
            .tab {{ padding: 10px; }}
        }}
    </style>
</head>
<body>
    <div class='container'>
        <h1>Azure FinOps Report</h1>
        <div class='sub'>Generated {escape(str(response.get("generatedAt")))} | Period {escape(str(response.get("period", {}).get("startDate")))} to {escape(str(response.get("period", {}).get("endDate")))}</div>
        <div class='sub'>Scope: Tenant {escape(str(response.get("scope", {}).get("tenantId", "auto")))} | Subscription {escape(str(response.get("scope", {}).get("subscriptionId", "auto")))}</div>
        <div class='scope-actions'>
            <button id='scope-switch-btn' class='scope-btn' type='button'>Switch Tenant / Subscription</button>
            <button id='scope-login-btn' class='scope-login' type='button'>Sign in with Microsoft</button>
            <button id='scope-logout-btn' class='scope-logout' type='button'>Sign out</button>
            <span id='scope-login-status' class='scope-status'>Not signed in</span>
            <span class='scope-help'>Uses your selected scope for this report request.</span>
        </div>
        <div class='artifacts'>
            <a class='artifact-link' href='{escape(str(response.get("artifactUrls", {}).get("pdf", "#")))}' target='_blank' rel='noopener'>Download PDF</a>
            <a class='artifact-link' href='{escape(str(response.get("artifactUrls", {}).get("powerpoint", "#")))}' target='_blank' rel='noopener'>Download PPTX</a>
            <a class='artifact-link' href='{escape(str(response.get("artifactUrls", {}).get("database", "#")))}' target='_blank' rel='noopener'>Download SQLite</a>
        </div>
        <dialog id='scope-dialog' class='scope-dialog'>
            <h3 style='margin:0;'>Run Against Another Tenant</h3>
            <p class='scope-help' style='margin:8px 0 0;'>Enter target tenant and subscription IDs, then click Sign in with Microsoft. Apply Scope uses that signed-in token.</p>
            <div class='scope-grid'>
                <label>Tenant ID
                    <input id='scope-tenant-id' type='text' value='{escape(str("" if response.get("scope", {}).get("tenantId") == "auto" else response.get("scope", {}).get("tenantId", "")))}' placeholder='e.g. e594a530-1ec9-4192-a8d4-a9111f8cffa7' />
                </label>
                <label>Subscription ID
                    <input id='scope-subscription-id' type='text' value='{escape(str("" if response.get("scope", {}).get("subscriptionId") == "auto" else response.get("scope", {}).get("subscriptionId", "")))}' placeholder='e.g. 5755893a-8056-4ba8-9916-1133c80a80f3' />
                </label>
            </div>
            <div class='scope-dialog-actions'>
                <button id='scope-cancel-btn' type='button'>Cancel</button>
                <button id='scope-apply-btn' class='apply' type='button'>Apply Scope</button>
            </div>
        </dialog>
        <div class='tabs-nav'>{''.join(nav_buttons)}</div>
        {''.join(sections)}
    </div>
    <script>
        function syncTableScrollbars() {{
            const sections = document.querySelectorAll('.tab');
            sections.forEach((section) => {{
                const topScroll = section.querySelector('.top-scroll');
                const topInner = section.querySelector('.top-scroll-inner');
                const tableWrap = section.querySelector('.table-wrap');
                const table = section.querySelector('table');
                if (!topScroll || !topInner || !tableWrap || !table) return;

                topInner.style.width = `${{table.scrollWidth}}px`;

                let syncing = false;
                topScroll.addEventListener('scroll', () => {{
                    if (syncing) return;
                    syncing = true;
                    tableWrap.scrollLeft = topScroll.scrollLeft;
                    syncing = false;
                }});

                tableWrap.addEventListener('scroll', () => {{
                    if (syncing) return;
                    syncing = true;
                    topScroll.scrollLeft = tableWrap.scrollLeft;
                    syncing = false;
                }});
            }});
        }}

        const buttons = document.querySelectorAll('.tab-btn');
        const tabsNodes = document.querySelectorAll('.tab');
        buttons.forEach((button) => {{
            button.addEventListener('click', () => {{
                buttons.forEach((b) => b.classList.remove('active'));
                tabsNodes.forEach((t) => t.classList.add('hidden'));
                button.classList.add('active');
                const target = document.getElementById(button.dataset.target);
                if (target) target.classList.remove('hidden');
                syncTableScrollbars();
            }});
        }});
        window.addEventListener('resize', syncTableScrollbars);
        syncTableScrollbars();

        const scopeDialog = document.getElementById('scope-dialog');
        const switchBtn = document.getElementById('scope-switch-btn');
        const loginBtn = document.getElementById('scope-login-btn');
        const logoutBtn = document.getElementById('scope-logout-btn');
        const loginStatus = document.getElementById('scope-login-status');
        const cancelBtn = document.getElementById('scope-cancel-btn');
        const applyBtn = document.getElementById('scope-apply-btn');
        const tenantInput = document.getElementById('scope-tenant-id');
        const subscriptionInput = document.getElementById('scope-subscription-id');
        const tokenStorageKey = 'finopsArmToken';
        const accountStorageKey = 'finopsArmAccount';
        const clientId = {json.dumps(MSAL_PUBLIC_CLIENT_ID)};

        function renderLoginStatus() {{
            const account = sessionStorage.getItem(accountStorageKey);
            if (account) {{
                loginStatus.textContent = `Signed in as ${{account}}`;
            }} else {{
                loginStatus.textContent = 'Not signed in';
            }}
        }}

        async function signInForArm() {{
            if (!clientId) {{
                alert('Interactive Azure sign-in is not configured. Set FINOPS_MSAL_CLIENT_ID to this tenant\'s Entra application client ID.');
                return;
            }}

            if (!window.msal) {{
                alert('Microsoft login library did not load. Refresh and try again.');
                return;
            }}

            const rawTenant = (tenantInput?.value || '').trim();
            const authorityTenant = rawTenant || 'organizations';
            const msalConfig = {{
                auth: {{
                    clientId,
                    authority: `https://login.microsoftonline.com/${{authorityTenant}}`,
                    redirectUri: window.location.origin + window.location.pathname,
                }}
            }};

            const app = new msal.PublicClientApplication(msalConfig);
            if (typeof app.initialize === 'function') {{
                await app.initialize();
            }}

            const loginRequest = {{
                scopes: ['https://management.azure.com/user_impersonation', 'openid', 'profile'],
                prompt: 'select_account',
            }};

            const loginResult = await app.loginPopup(loginRequest);
            const account = loginResult.account;
            if (!account) {{
                throw new Error('No account returned from Microsoft login');
            }}

            const tokenResult = await app.acquireTokenSilent({{
                account,
                scopes: ['https://management.azure.com/user_impersonation'],
            }});

            sessionStorage.setItem(tokenStorageKey, tokenResult.accessToken || '');
            sessionStorage.setItem(accountStorageKey, account.username || account.homeAccountId || 'signed-in user');
            renderLoginStatus();
        }}

        function signOutArm() {{
            sessionStorage.removeItem(tokenStorageKey);
            sessionStorage.removeItem(accountStorageKey);
            renderLoginStatus();
        }}

        renderLoginStatus();

        if (switchBtn && scopeDialog) {{
            switchBtn.addEventListener('click', () => {{
                if (typeof scopeDialog.showModal === 'function') {{
                    scopeDialog.showModal();
                }} else {{
                    scopeDialog.setAttribute('open', 'open');
                }}
            }});
        }}

        if (cancelBtn && scopeDialog) {{
            cancelBtn.addEventListener('click', () => {{
                if (typeof scopeDialog.close === 'function') {{
                    scopeDialog.close();
                }} else {{
                    scopeDialog.removeAttribute('open');
                }}
            }});
        }}

        if (loginBtn) {{
            loginBtn.addEventListener('click', async () => {{
                try {{
                    await signInForArm();
                }} catch (err) {{
                    const message = err && err.message ? err.message : String(err);
                    alert(`Sign-in failed: ${{message}}`);
                }}
            }});
        }}

        if (logoutBtn) {{
            logoutBtn.addEventListener('click', () => {{
                signOutArm();
            }});
        }}

        if (applyBtn) {{
            applyBtn.addEventListener('click', async () => {{
                const url = new URL(window.location.href);
                const tenantId = (tenantInput?.value || '').trim();
                const subscriptionId = (subscriptionInput?.value || '').trim();

                if (tenantId) {{
                    url.searchParams.set('tenantId', tenantId);
                }} else {{
                    url.searchParams.delete('tenantId');
                }}

                if (subscriptionId) {{
                    url.searchParams.set('subscriptionId', subscriptionId);
                }} else {{
                    url.searchParams.delete('subscriptionId');
                }}

                url.searchParams.set('format', 'html');
                const armToken = sessionStorage.getItem(tokenStorageKey) || '';

                if (!armToken) {{
                    alert('Sign in first so the report can run with your tenant/subscription access.');
                    return;
                }}

                const response = await fetch(url.toString(), {{
                    method: 'GET',
                    headers: {{
                        'Accept': 'text/html',
                        'Authorization': `Bearer ${{armToken}}`,
                    }}
                }});

                const html = await response.text();
                document.open();
                document.write(html);
                document.close();
            }});
        }}
    </script>
</body>
</html>
"""


def _run_finops_report_request(req: func.HttpRequest) -> func.HttpResponse:
    try:
        page = _parse_positive_int(req.params.get("page"), 1)
        # Enforce capped payload size for faster responses.
        page_size = min(_parse_positive_int(req.params.get("pageSize"), 100), 250)
        target_tenant_id = (req.params.get("tenantId") or "").strip() or None
        target_subscription_id = (req.params.get("subscriptionId") or "").strip() or None
        auth_header = (req.headers.get("Authorization") or "").strip()
        provided_arm_token = None
        if auth_header.lower().startswith("bearer "):
            provided_arm_token = auth_header[7:].strip() or None
        if not provided_arm_token:
            provided_arm_token = (req.headers.get("x-arm-access-token") or "").strip() or None

        snapshot = collect_report_data(
            target_tenant_id=target_tenant_id,
            target_subscription_id=target_subscription_id,
            arm_access_token=provided_arm_token,
        )
        artifacts = generate_artifacts(snapshot)
        artifact_urls = _build_artifact_urls(req)
        response = build_tabular_output(snapshot, artifacts, artifact_urls, page, page_size)
        format_param = (req.params.get("format") or "").lower()
        accept_header = (req.headers.get("Accept") or "").lower()
        wants_html = format_param == "html" or ("text/html" in accept_header and format_param != "json")
        if wants_html:
            return func.HttpResponse(render_tabular_html(response), mimetype="text/html")
        return func.HttpResponse(json.dumps(response), mimetype="application/json")
    except Exception as exc:
        logging.exception("FinOps report generation failed")
        return func.HttpResponse(f"Report generation failed: {exc}", status_code=500)


@app.function_name(name="run_finops_report")
@app.route(route="run_finops_report", methods=["GET", "POST"], auth_level=HTTP_ROUTE_AUTH_LEVEL)
def run_finops_report(req: func.HttpRequest) -> func.HttpResponse:
    return _run_finops_report_request(req)


@app.function_name(name="finops_health")
@app.route(route="finops_health", methods=["GET"], auth_level=HTTP_ROUTE_AUTH_LEVEL)
def finops_health(req: func.HttpRequest) -> func.HttpResponse:
    payload = {
        "status": "ok",
        "azureCli": "unavailable",
        "managedIdentity": {
            "status": "unknown",
            "tokenAcquired": False,
        },
        "subscriptionHint": _CURRENT_SUBSCRIPTION_ID or _subscription_id_from_owner_name() or "",
        "artifactDir": str(ARTIFACT_DIR),
        "reportDbPath": str(REPORT_DB_PATH),
    }

    try:
        _get_arm_access_token()
        payload["managedIdentity"]["status"] = "ok"
        payload["managedIdentity"]["tokenAcquired"] = True
    except Exception as exc:
        payload["status"] = "degraded"
        payload["managedIdentity"]["status"] = "error"
        payload["managedIdentity"]["error"] = str(exc)

    try:
        az_path = _resolve_az_executable()
        payload["azureCli"] = "available"
        payload["azureCliPath"] = az_path
    except RuntimeError as exc:
        payload["message"] = str(exc)

    return func.HttpResponse(json.dumps(payload), mimetype="application/json")


@app.function_name(name="finops_artifact")
@app.route(route="finops_artifact", methods=["GET"], auth_level=HTTP_ROUTE_AUTH_LEVEL)
def finops_artifact(req: func.HttpRequest) -> func.HttpResponse:
    artifact_name = (req.params.get("name") or req.params.get("type") or "").strip().lower()
    artifact_files = {
        "database": ARTIFACT_DIR / "finops_reports.db",
        "pdf": ARTIFACT_DIR / "finops_summary.pdf",
        "powerpoint": ARTIFACT_DIR / "finops_summary.pptx",
    }
    target = artifact_files.get(artifact_name)
    if target is None:
        return func.HttpResponse("Unknown artifact name", status_code=400)
    if not target.exists() or not target.is_file():
        return func.HttpResponse("Artifact not found. Run run_finops_report first.", status_code=404)

    mimetype = {
        "database": "application/octet-stream",
        "pdf": "application/pdf",
        "powerpoint": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }[artifact_name]

    with open(target, "rb") as handle:
        payload = handle.read()

    return func.HttpResponse(
        body=payload,
        mimetype=mimetype,
        headers={"Content-Disposition": f'attachment; filename="{target.name}"'},
    )


@app.function_name(name="scheduled_finops_report")
@app.schedule(schedule="0 0 6 * * *", arg_name="mytimer", run_on_startup=False)
def scheduled_finops_report(mytimer: func.TimerRequest) -> None:
    if mytimer.past_due:
        logging.info("Timer is past due")
    snapshot = collect_report_data()
    artifacts = generate_artifacts(snapshot)
    logging.info("Scheduled FinOps run completed: %s", artifacts)
